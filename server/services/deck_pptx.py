"""Native .pptx generation from a validated deck spec (python-pptx).

Safe by construction, exactly like render_chart: the model supplies a normalized DATA spec
(slides + content + a closed set of layout intents), never code. We build the PowerPoint from
that spec. No shell, no network, no untrusted-file parsing (v1 ships clean themed decks;
template-following is deferred to v2). Output is a real, natively-editable .pptx (shapes + text
frames + speaker notes) — not slide images.

Inspired by hugohe3/ppt-master (MIT) for the "native editable shapes, not images" approach.
The design system (mandatory 5-role color theme + strict typography scale + per-page-type
composition rules) is a port of the SPEC from MiniMax-AI/skills pptx-generator (MIT) — no code
was vendored. See THIRD_PARTY_NOTICES.md.
"""
from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

LAYOUTS = {"title", "section", "bullets", "two-column", "quote", "big-number",
           "table", "chart", "kpi"}

# 16:9 canvas.
_W = Inches(13.333)
_H = Inches(7.5)
_MARGIN = Inches(0.6)
_CONTENT_W = _W - _MARGIN * 2

# Typeface pair applied to EVERY run. python-pptx's font.name only sets the latin
# (western-script) typeface — CJK text silently falls back to the viewer default (Calibri-ish,
# looks cheap) unless an <a:ea> element is ALSO written into the run properties. PingFang SC is
# the macOS system CJK face; PowerPoint elsewhere falls back gracefully (e.g. Microsoft YaHei).
_LATIN_FONT = "Helvetica Neue"
_EA_FONT = "PingFang SC"

# Typography scale (pt) — consistent across every layout.
_T_TITLE = 40    # title-slide headline / section-slide oversize title
_T_HEADING = 32  # content-slide headings
_T_QUOTE = 28    # quote body
_T_BODY = 18     # bullets / subtitles / labels
_T_SUB = 16      # sub-bullets
_T_CAPTION = 12  # page badge / attributions / column headers

# Design system: every theme is a mandatory set of 5 color roles (hex, RGBColor-ready).
#   primary   — dominant ink: body text + solid bands (high contrast against bg)
#   secondary — supporting text: subtitles, attributions, sub-bullets, column headers
#   accent    — highlights ONLY: underline bars, badges, big section numbers;
#               never body text on a light background
#   light     — text placed ON a primary-filled band
#   bg        — slide background
THEMES: dict[str, dict[str, str]] = {
    # The house look — the user's own WCAG-AA calibrated "Ember" system: wheat paper,
    # warm black ink, rust-orange accent (4.6:1 on paper; _on_accent resolves to white
    # text on this dark accent automatically).
    "ember": {"primary": "1A1410", "secondary": "6B5E52", "accent": "D94420",
              "light": "FAF5EE", "bg": "F2EBE0"},
    # Ember dark variant: warm black bg, paper text, brighter on-dark accent.
    "ember-dark": {"primary": "F2EBE0", "secondary": "D0A890", "accent": "F06A20",
                   "light": "262018", "bg": "1A1410"},
    # Previous house look: warm paper, near-black ink, chartreuse highlighter.
    "ink": {"primary": "1E1E1E", "secondary": "54665A", "accent": "EEFF53",
            "light": "FFFFFF", "bg": "F6F4F2"},
    # Dark: near-black blue, light text, cyan accent.
    "midnight": {"primary": "E8EEF4", "secondary": "9AA8B5", "accent": "22D3EE",
                 "light": "101418", "bg": "101418"},
    # Clean corporate: white, deep blue ink, sky accent.
    "azure": {"primary": "14337A", "secondary": "4A6591", "accent": "38BDF8",
              "light": "FFFFFF", "bg": "FFFFFF"},
    # Warm editorial: cream, brown/rust ink, orange accent.
    "terra": {"primary": "7C3F21", "secondary": "9A6B4F", "accent": "E8852B",
              "light": "FFF8EE", "bg": "FAF3E7"},
}
DEFAULT_THEME = "ember"


def _resolve_theme(name: object) -> dict[str, RGBColor]:
    """Resolve a preset name to RGBColor roles; anything unknown (or non-string, e.g. the
    legacy {'accent': '#hex'} dict) falls back to the default preset — never crashes."""
    key = name.strip().lower() if isinstance(name, str) else ""
    hexes = THEMES.get(key, THEMES[DEFAULT_THEME])
    t = {role: RGBColor.from_string(h) for role, h in hexes.items()}
    # Text sitting on the accent: pick by the ACCENT's own luminance — a bright accent
    # (lime, cyan) needs the darker of primary/light on it, a dark accent (ember's rust
    # #D94420) needs the lighter one (white-ish). The old darker-always rule assumed
    # bright accents and would put near-black on rust.
    def _luma(c: RGBColor) -> float:
        return 0.299 * c[0] + 0.587 * c[1] + 0.114 * c[2]
    pick = min if _luma(t["accent"]) >= 128 else max
    t["on_accent"] = pick((t["primary"], t["light"]), key=_luma)
    return t


def _rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(slide, left, top, width, height, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    return box, tf, p


def _set_font(run, *, size=None, color=None, bold=None, italic=None):
    """Apply size/weight/color AND both typefaces to a run. font.name only writes the
    <a:latin> typeface, so the east-asian one is set via a raw <a:ea> element — without it,
    Chinese text renders in the viewer's default fallback font."""
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if color is not None:
        f.color.rgb = color
    f.name = _LATIN_FONT  # writes <a:latin> (and guarantees rPr exists)
    rPr = f._rPr
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.find(qn("a:latin")).addnext(ea)  # schema order: a:ea directly follows a:latin
    ea.set("typeface", _EA_FONT)
    return run


def _run(paragraph, text, *, size, color, bold=False, italic=False):
    r = paragraph.add_run()
    r.text = str(text)
    return _set_font(r, size=size, color=color, bold=bold, italic=italic)


def _accent_underline(slide, t, *, top):
    """Short accent underline bar under a heading (~1.2in x 4pt)."""
    return _rect(slide, _MARGIN, top, Inches(1.2), Pt(4), t["accent"])


def _page_badge(slide, number, t):
    """Small accent-filled rectangle with the page number, bottom-right."""
    w, h = Inches(0.55), Inches(0.32)
    shp = _rect(slide, _W - Inches(0.85), _H - Inches(0.57), w, h, t["accent"])
    tf = shp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, attr, 0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, str(number), size=_T_CAPTION, color=t["on_accent"], bold=True)


def _bullet_list(slide, items, t, *, left, top, width, height):
    """Bullets in primary; items prefixed '- ' or indented render as sub-bullets in secondary."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text = str(item)
        sub = text.startswith(("- ", "  "))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        if sub:
            _run(p, f"     –  {text.lstrip('- ').strip()}", size=_T_SUB, color=t["secondary"])
        else:
            _run(p, f"•  {text}", size=_T_BODY, color=t["primary"])
    return box


def _heading(slide, title, t):
    """Content-slide heading: 32pt primary + short accent underline bar."""
    _, _, p = _text(slide, _MARGIN, Inches(0.55), _CONTENT_W, Inches(0.9))
    _run(p, title or "", size=_T_HEADING, color=t["primary"], bold=True)
    _accent_underline(slide, t, top=Inches(1.42))


def _title_slide(slide, s, t):
    # Full-bleed primary band across the lower-middle third; title in light on the band.
    band_top, band_h = Inches(3.1), Inches(2.7)
    _rect(slide, 0, band_top, _W, band_h, t["primary"])
    # Small accent eyebrow bar above the title.
    _rect(slide, _MARGIN, band_top + Inches(0.45), Inches(1.4), Inches(0.09), t["accent"])
    _, _, p = _text(slide, _MARGIN, band_top + Inches(0.75), _CONTENT_W, Inches(1.5))
    _run(p, s.get("title") or "", size=_T_TITLE, color=t["light"], bold=True)
    if s.get("subtitle"):
        _, _, sp = _text(slide, _MARGIN, band_top + band_h + Inches(0.25), _CONTENT_W, Inches(0.8))
        _run(sp, s["subtitle"], size=_T_BODY, color=t["secondary"])


def _section_slide(slide, s, t, number):
    # Big section number in accent, oversize title in primary, thin accent rule between.
    _, _, np = _text(slide, _MARGIN, Inches(1.0), _CONTENT_W, Inches(2.2))
    _run(np, f"{number:02d}", size=96, color=t["accent"], bold=True)
    _rect(slide, _MARGIN, Inches(3.35), Inches(2.4), Pt(3), t["accent"])
    _, _, p = _text(slide, _MARGIN, Inches(3.7), _CONTENT_W, Inches(1.6))
    _run(p, s.get("title") or "", size=_T_TITLE, color=t["primary"], bold=True)


_MAX_BULLETS_SHOWN = 6  # per list; overflow is truncated on-slide, full list goes to notes


def _cap_bullets(items, label):
    """Cap a bullet list at _MAX_BULLETS_SHOWN. Returns (shown_items, extra_note_or_None);
    the note carries the FULL list so nothing is lost — it lands in the speaker notes."""
    items = [str(x) for x in items]
    if len(items) <= _MAX_BULLETS_SHOWN:
        return items, None
    note = (f"[{label} truncated to {_MAX_BULLETS_SHOWN} on the slide] Full list:\n"
            + "\n".join(f"• {x}" for x in items))
    return items[:_MAX_BULLETS_SHOWN], note


def _bullets_slide(slide, s, t):
    _heading(slide, s.get("title"), t)
    shown, extra = _cap_bullets(s.get("bullets") or [], "bullets")
    _bullet_list(slide, shown, t,
                 left=_MARGIN, top=Inches(1.95), width=_CONTENT_W, height=Inches(4.7))
    return extra


def _two_col_slide(slide, s, t):
    _heading(slide, s.get("title"), t)
    gap = Inches(0.7)
    col_w = (_CONTENT_W - gap) / 2
    right_x = _MARGIN + col_w + gap
    # Hairline vertical divider between the columns.
    _rect(slide, _MARGIN + col_w + gap / 2, Inches(1.95), Pt(1), Inches(4.6), t["secondary"])
    body_top = Inches(1.95)
    headers = (s.get("left_title"), s.get("right_title"))
    if any(headers):
        body_top = Inches(2.55)
        for x, header in ((_MARGIN, headers[0]), (right_x, headers[1])):
            if header:
                _, _, hp = _text(slide, x, Inches(1.95), col_w, Inches(0.4))
                _run(hp, str(header).upper(), size=_T_CAPTION + 2, color=t["secondary"], bold=True)
            _rect(slide, x, Inches(2.4), col_w, Pt(1), t["secondary"])  # hairline under header row
    left_shown, ln = _cap_bullets(s.get("left") or [], "left column")
    right_shown, rn = _cap_bullets(s.get("right") or [], "right column")
    _bullet_list(slide, left_shown, t,
                 left=_MARGIN, top=body_top, width=col_w, height=Inches(4.2))
    _bullet_list(slide, right_shown, t,
                 left=right_x, top=body_top, width=col_w, height=Inches(4.2))
    extra = "\n\n".join(x for x in (ln, rn) if x)
    return extra or None


def _quote_slide(slide, s, t):
    # Oversized accent quotation mark, quote in primary italic, attribution in secondary.
    _, _, qp = _text(slide, _MARGIN, Inches(0.8), Inches(2.5), Inches(1.8))
    _run(qp, "“", size=120, color=t["accent"], bold=True)
    _, _, p = _text(slide, _MARGIN + Inches(0.3), Inches(2.4), _CONTENT_W - Inches(0.6),
                    Inches(2.6), anchor=MSO_ANCHOR.MIDDLE)
    _run(p, s.get("text") or "", size=_T_QUOTE, color=t["primary"], bold=True, italic=True)
    if s.get("attribution"):
        _, _, ap = _text(slide, _MARGIN + Inches(0.3), Inches(5.3), _CONTENT_W - Inches(0.6),
                         Inches(0.7))
        _run(ap, f"— {s['attribution']}", size=_T_BODY, color=t["secondary"])


def _big_number_slide(slide, s, t):
    # The number huge in primary with an accent underline block, caption in secondary.
    _, _, p = _text(slide, _MARGIN, Inches(1.9), _CONTENT_W, Inches(2.2),
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _run(p, s.get("value") or "", size=96, color=t["primary"], bold=True)
    _rect(slide, (_W - Inches(2.0)) / 2, Inches(4.35), Inches(2.0), Inches(0.12), t["accent"])
    if s.get("label"):
        _, _, lp = _text(slide, _MARGIN, Inches(4.75), _CONTENT_W, Inches(0.9),
                         align=PP_ALIGN.CENTER)
        _run(lp, s["label"], size=_T_BODY, color=t["secondary"])


# ---- data layouts (table / chart / kpi) — numbers deserve better than bullet walls ----

_TABLE_MAX_ROWS = 7  # body rows shown; extra rows are truncated into the speaker notes


def _cell_text(cell, text, *, size, color, bold=False):
    cell.margin_left = cell.margin_right = Inches(0.1)
    cell.margin_top = cell.margin_bottom = Inches(0.05)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    _run(tf.paragraphs[0], text, size=size, color=color, bold=bold)


def _cell_bottom_accent(cell, color):
    """Accent bottom edge on a header cell — python-pptx has no cell-border API, so write
    the <a:lnB> element into tcPr directly (borders precede fill in the tcPr sequence)."""
    tcPr = cell._tc.get_or_add_tcPr()
    lnB = tcPr.find(qn("a:lnB"))
    if lnB is None:
        lnB = tcPr.makeelement(qn("a:lnB"), {})
        tcPr.insert(0, lnB)
    lnB.set("w", "28575")  # 2.25pt (EMU)
    fill = lnB.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"), {"val": str(color)})
    fill.append(clr)
    lnB.append(fill)


def _table_slide(slide, s, t):
    _heading(slide, s.get("title"), t)
    headers = [str(h) for h in (s.get("headers") or [])]
    all_rows = [[str(c) for c in row] for row in (s.get("rows") or [])]
    rows = all_rows[:_TABLE_MAX_ROWS]
    top = Inches(1.95)
    gf = slide.shapes.add_table(len(rows) + 1, len(headers), _MARGIN, top,
                                _CONTENT_W, Inches(0.55) * (len(rows) + 1))
    table = gf.table
    table.horz_banding = False  # we zebra explicitly with theme colors
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = t["primary"]
        _cell_bottom_accent(cell, t["accent"])
        _cell_text(cell, h, size=13, color=t["light"], bold=True)
    for i, row in enumerate(rows):
        fill = t["light"] if i % 2 else t["bg"]  # zebra
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            _cell_text(cell, val, size=13, color=t["primary"])
    if len(all_rows) > _TABLE_MAX_ROWS:
        hidden = all_rows[_TABLE_MAX_ROWS:]
        return (f"[table truncated — showing {_TABLE_MAX_ROWS} of {len(all_rows)} rows] "
                "Remaining rows:\n" + "\n".join(" | ".join(r) for r in hidden))
    return None


_CHART_TYPES = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}


def _chart_slide(slide, s, t):
    """Native, editable pptx chart (not an image). Data problems raise ValueError — the
    executor surfaces the message to the model so it can fix the spec and retry."""
    ctype = s.get("chart_type")
    if ctype not in _CHART_TYPES:
        raise ValueError(f"chart_type must be one of: {', '.join(sorted(_CHART_TYPES))}")
    cats = [str(c) for c in (s.get("categories") or [])]
    series = s.get("series") or []
    if not cats or not series:
        raise ValueError("chart needs non-empty 'categories' and 'series'")
    data = CategoryChartData()
    data.categories = cats
    for ser in series:
        vals = [float(v) for v in (ser.get("values") or [])]
        if len(vals) != len(cats):
            raise ValueError(f"series {ser.get('name')!r} has {len(vals)} values but there are "
                             f"{len(cats)} categories — lengths must match")
        data.add_series(str(ser.get("name") or ""), vals)
    _heading(slide, s.get("title"), t)
    gf = slide.shapes.add_chart(_CHART_TYPES[ctype], _MARGIN, Inches(1.95),
                                _CONTENT_W, Inches(4.8), data)
    chart = gf.chart
    chart.has_title = False  # the styled slide heading is the title
    chart.has_legend = len(series) > 1  # legend only earns its space with multiple series
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    # Theme the first series where python-pptx allows (pie keeps auto-varied slice colors).
    first = chart.plots[0].series[0]
    if ctype == "line":
        first.format.line.color.rgb = t["accent"]
        first.format.line.width = Pt(2.5)
    elif ctype != "pie":
        first.format.fill.solid()
        first.format.fill.fore_color.rgb = t["accent"]


def _kpi_slide(slide, s, t):
    # Row of stat blocks: light rounded card + accent top bar, value huge, label quiet.
    _heading(slide, s.get("title"), t)
    items = list(s.get("items") or [])[:4]
    n = max(len(items), 1)
    gap = Inches(0.4)
    w = (_CONTENT_W - gap * (n - 1)) / n
    top, h = Inches(2.35), Inches(2.7)
    for i, item in enumerate(items):
        x = _MARGIN + (w + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = t["light"]
        card.line.color.rgb = t["secondary"]  # hairline keeps cards visible on same-color bgs
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        _rect(slide, x + Inches(0.35), top, w - Inches(0.7), Inches(0.09), t["accent"])
        _, _, vp = _text(slide, x + Inches(0.1), top + Inches(0.5), w - Inches(0.2), Inches(1.2),
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _run(vp, item.get("value") or "", size=_T_TITLE, color=t["primary"], bold=True)
        _, _, lp = _text(slide, x + Inches(0.15), top + Inches(1.8), w - Inches(0.3), Inches(0.7),
                         align=PP_ALIGN.CENTER)
        _run(lp, item.get("label") or "", size=_T_CAPTION, color=t["secondary"])


_RENDERERS = {
    "title": _title_slide, "section": _section_slide, "bullets": _bullets_slide,
    "two-column": _two_col_slide, "quote": _quote_slide, "big-number": _big_number_slide,
    "table": _table_slide, "chart": _chart_slide, "kpi": _kpi_slide,
}


def build_deck(spec: dict) -> bytes:
    """Build a .pptx from a validated deck spec and return its bytes. Assumes the spec has
    already passed DeckExecutor validation (layouts in LAYOUTS, bounded counts).
    spec['theme'] is an optional preset name from THEMES; unknown/absent -> default."""
    t = _resolve_theme(spec.get("theme"))

    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H
    blank = prs.slide_layouts[6]  # blank — we place every shape ourselves for full theme control

    section_no = 0
    for idx, s in enumerate(spec["slides"], start=1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = t["bg"]
        layout = s["layout"]
        if layout == "section":
            section_no += 1
            extra_notes = _section_slide(slide, s, t, section_no)
        else:
            extra_notes = _RENDERERS[layout](slide, s, t)
        if layout != "title":
            _page_badge(slide, idx, t)
        # Speaker notes = the author's notes + anything a renderer truncated off the slide.
        parts = [str(s["notes"])] if s.get("notes") else []
        if extra_notes:
            parts.append(extra_notes)
        if parts:
            slide.notes_slide.notes_text_frame.text = "\n\n".join(parts)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
