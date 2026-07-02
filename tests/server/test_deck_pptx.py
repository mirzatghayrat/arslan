import base64
from io import BytesIO

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

from server.registry.executors import DeckExecutor
from server.services.deck_pptx import DEFAULT_THEME, THEMES

pytestmark = pytest.mark.asyncio

VALID = {
    "title": "AI 2026",
    "slides": [
        {"layout": "title", "title": "AI 2026", "subtitle": "A field briefing"},
        {"layout": "section", "title": "The Landscape"},
        {"layout": "bullets", "title": "Trends", "bullets": ["Agents", "Cheaper inference"], "notes": "speak"},
        {"layout": "two-column", "title": "Bull / Bear", "left": ["up"], "right": ["moats"],
         "left_title": "Bull", "right_title": "Bear"},
        {"layout": "quote", "text": "The future is here.", "attribution": "Gibson"},
        {"layout": "big-number", "value": "92%", "label": "devs using AI"},
    ],
}

DATA_SLIDES = [
    {"layout": "table", "title": "模型对比",
     "headers": ["Model", "Score", "Cost"],
     "rows": [["A", "92", "$1"], ["B", "88", "$2"], ["C", "85", "$3"]]},
    {"layout": "chart", "title": "季度收入", "chart_type": "column",
     "categories": ["Q1", "Q2", "Q3"], "series": [{"name": "收入", "values": [1.2, 3.4, 5.6]}]},
    {"layout": "chart", "title": "Share", "chart_type": "pie",
     "categories": ["A", "B", "C"], "series": [{"name": "share", "values": [50, 30, 20]}]},
    {"layout": "kpi", "title": "At a glance",
     "items": [{"value": "92%", "label": "采用率"}, {"value": "3.4x", "label": "增速"},
               {"value": "$12M", "label": "ARR"}]},
]


def _open(out) -> Presentation:
    assert out["ok"] is True, out.get("error")
    return Presentation(BytesIO(base64.b64decode(out["artifact"]["bytes_b64"])))


def _solid_fills(slide) -> set[RGBColor]:
    """All solid fill colors on a slide's shapes."""
    fills = set()
    for shp in slide.shapes:
        try:
            if shp.fill.type == MSO_FILL.SOLID:
                fills.add(shp.fill.fore_color.rgb)
        except (AttributeError, TypeError):
            continue
    return fills


def _theme_rgb(name: str, role: str) -> RGBColor:
    return RGBColor.from_string(THEMES[name][role])


async def test_render_deck_produces_valid_editable_pptx():
    out = await DeckExecutor().execute(VALID)
    assert out["ok"] is True
    art = out["artifact"]
    assert art["kind"] == "pptx" and art["filename"].endswith(".pptx") and art["slides"] == 6
    # the artifact is a real, re-openable PowerPoint with native shapes (not an image)
    prs = Presentation(BytesIO(base64.b64decode(art["bytes_b64"])))
    assert len(prs.slides._sldIdLst) == 6
    # speaker notes survived, and every slide carries native shapes
    assert prs.slides[2].notes_slide.notes_text_frame.text == "speak"
    assert all(len(s.shapes) >= 1 for s in prs.slides)


async def test_default_theme_ember_applied():
    prs = _open(await DeckExecutor().execute(VALID))
    assert DEFAULT_THEME == "ember"
    # every slide gets the theme background fill
    for slide in prs.slides:
        assert slide.background.fill.fore_color.rgb == _theme_rgb("ember", "bg")
    # title slide carries the primary-color band + accent eyebrow bar
    title_fills = _solid_fills(prs.slides[0])
    assert _theme_rgb("ember", "primary") in title_fills
    assert _theme_rgb("ember", "accent") in title_fills
    # bullets slide carries the accent underline bar (and accent page badge)
    assert _theme_rgb("ember", "accent") in _solid_fills(prs.slides[2])


async def test_explicit_theme_applied():
    spec = dict(VALID, theme="midnight")
    prs = _open(await DeckExecutor().execute(spec))
    assert prs.slides[0].background.fill.fore_color.rgb == _theme_rgb("midnight", "bg")
    title_fills = _solid_fills(prs.slides[0])
    assert _theme_rgb("midnight", "primary") in title_fills
    assert _theme_rgb("midnight", "accent") in _solid_fills(prs.slides[2])


async def test_unknown_theme_falls_back_to_default():
    prs = _open(await DeckExecutor().execute(dict(VALID, theme="vaporwave")))
    assert prs.slides[0].background.fill.fore_color.rgb == _theme_rgb(DEFAULT_THEME, "bg")


async def test_legacy_dict_theme_falls_back_to_default():
    # older callers passed theme={'accent': '#hex'}; must not crash, must fall back
    prs = _open(await DeckExecutor().execute(dict(VALID, theme={"accent": "#2b6ee8"})))
    assert prs.slides[0].background.fill.fore_color.rgb == _theme_rgb(DEFAULT_THEME, "bg")


async def test_all_presets_generate_without_error():
    for name in THEMES:
        prs = _open(await DeckExecutor().execute(dict(VALID, theme=name)))
        assert len(prs.slides._sldIdLst) == 6
        assert prs.slides[0].background.fill.fore_color.rgb == _theme_rgb(name, "bg")


async def test_page_badge_on_all_slides_except_title():
    prs = _open(await DeckExecutor().execute(VALID))
    accent = _theme_rgb("ember", "accent")
    for i, slide in enumerate(prs.slides):
        badge_texts = {
            shp.text_frame.text for shp in slide.shapes
            if shp.has_text_frame and shp.text_frame.text.strip().isdigit()
        }
        if i == 0:
            assert str(i + 1) not in badge_texts  # no page number on the title slide
        else:
            assert str(i + 1) in badge_texts
            assert accent in _solid_fills(slide)


async def test_empty_slides_rejected():
    out = await DeckExecutor().execute({"slides": []})
    assert out["ok"] is False and "non-empty" in out["error"]


async def test_invalid_layout_rejected():
    out = await DeckExecutor().execute({"slides": [{"layout": "carousel", "title": "x"}]})
    assert out["ok"] is False and "invalid layout" in out["error"]


async def test_bullets_overflow_rejected():
    out = await DeckExecutor().execute(
        {"slides": [{"layout": "bullets", "title": "x", "bullets": [str(i) for i in range(99)]}]}
    )
    assert out["ok"] is False and "bullets" in out["error"]


async def test_too_many_slides_rejected():
    out = await DeckExecutor().execute({"slides": [{"layout": "section", "title": "s"}] * 99})
    assert out["ok"] is False and "too many" in out["error"]


async def test_table_layout_native_table_with_themed_header():
    out = await DeckExecutor().execute({"slides": [DATA_SLIDES[0]]})
    prs = _open(out)
    tables = [shp.table for shp in prs.slides[0].shapes if shp.has_table]
    assert len(tables) == 1  # a real GraphicFrame table, not a textbox of numbers
    tbl = tables[0]
    assert [tbl.cell(0, j).text for j in range(3)] == ["Model", "Score", "Cost"]
    assert tbl.cell(1, 0).text == "A" and tbl.cell(3, 2).text == "$3"
    # header row is filled with the theme primary
    assert tbl.cell(0, 0).fill.fore_color.rgb == _theme_rgb(DEFAULT_THEME, "primary")


async def test_table_rows_truncated_to_seven_with_full_data_in_notes():
    rows = [[f"r{i}", str(i)] for i in range(12)]
    out = await DeckExecutor().execute(
        {"slides": [{"layout": "table", "title": "t", "headers": ["k", "v"], "rows": rows}]})
    prs = _open(out)
    tbl = next(shp.table for shp in prs.slides[0].shapes if shp.has_table)
    assert len(tbl.rows) == 8  # header + 7 body rows
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "r7" in notes and "r11" in notes  # overflow preserved in speaker notes


async def test_table_invalid_spec_rejected():
    out = await DeckExecutor().execute(
        {"slides": [{"layout": "table", "title": "t", "headers": ["k", "v"],
                     "rows": [["only-one-cell"]]}]})
    assert out["ok"] is False and "headers" in out["error"]


async def test_chart_layout_native_editable_chart_column_and_pie():
    out = await DeckExecutor().execute({"slides": [DATA_SLIDES[1], DATA_SLIDES[2]]})
    prs = _open(out)
    col = next(shp.chart for shp in prs.slides[0].shapes if shp.has_chart)
    assert col.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
    assert list(col.plots[0].categories) == ["Q1", "Q2", "Q3"]
    assert list(col.plots[0].series[0].values) == [1.2, 3.4, 5.6]
    pie = next(shp.chart for shp in prs.slides[1].shapes if shp.has_chart)
    assert pie.chart_type == XL_CHART_TYPE.PIE


async def test_chart_mismatched_values_rejected_with_clear_error():
    bad = {"layout": "chart", "title": "x", "chart_type": "bar",
           "categories": ["a", "b"], "series": [{"name": "s", "values": [1]}]}
    out = await DeckExecutor().execute({"slides": [bad]})
    assert out["ok"] is False and "match" in out["error"]
    out = await DeckExecutor().execute(
        {"slides": [{"layout": "chart", "title": "x", "chart_type": "donut",
                     "categories": ["a"], "series": [{"values": [1]}]}]})
    assert out["ok"] is False and "chart_type" in out["error"]


async def test_kpi_layout_renders_stat_blocks():
    out = await DeckExecutor().execute({"slides": [DATA_SLIDES[3]]})
    prs = _open(out)
    cards = [shp for shp in prs.slides[0].shapes
             if shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and shp.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE]
    assert len(cards) == 3  # one rounded stat card per item
    texts = {shp.text_frame.text for shp in prs.slides[0].shapes if shp.has_text_frame}
    assert {"92%", "3.4x", "$12M", "采用率"} <= texts


async def test_kpi_item_count_enforced():
    for items in ([{"value": "1", "label": "x"}],
                  [{"value": str(i), "label": "x"} for i in range(5)]):
        out = await DeckExecutor().execute(
            {"slides": [{"layout": "kpi", "title": "t", "items": items}]})
        assert out["ok"] is False and "2-4" in out["error"]


async def test_east_asian_font_set_on_runs():
    prs = _open(await DeckExecutor().execute(VALID))
    box = next(shp for shp in prs.slides[0].shapes
               if shp.has_text_frame and "AI 2026" in shp.text_frame.text)
    run = box.text_frame.paragraphs[0].runs[0]
    xml = run._r.xml
    # latin via font.name AND a raw <a:ea> element — without the latter CJK falls back to Calibri
    assert run.font.name == "Helvetica Neue"
    assert "<a:ea" in xml and 'typeface="PingFang SC"' in xml


async def test_bullets_truncated_to_six_with_full_list_in_notes():
    bullets = [f"point {i}" for i in range(9)]
    out = await DeckExecutor().execute(
        {"slides": [{"layout": "bullets", "title": "t", "bullets": bullets, "notes": "keep me"}]})
    prs = _open(out)
    body = next(shp for shp in prs.slides[0].shapes
                if shp.has_text_frame and "point 0" in shp.text_frame.text)
    assert "point 5" in body.text_frame.text and "point 6" not in body.text_frame.text
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "keep me" in notes and "point 8" in notes  # author notes kept + full list appended


async def test_all_themes_generate_with_data_layouts():
    spec = {"slides": VALID["slides"] + DATA_SLIDES}
    for name in THEMES:
        prs = _open(await DeckExecutor().execute(dict(spec, theme=name)))
        assert len(prs.slides._sldIdLst) == len(spec["slides"])
        assert prs.slides[0].background.fill.fore_color.rgb == _theme_rgb(name, "bg")


async def test_deck_tool_is_registered_and_resolvable():
    from server.registry.executors import EXECUTORS, resolve_executor
    assert "render_deck" in EXECUTORS
    assert await resolve_executor("render_deck") is not None
